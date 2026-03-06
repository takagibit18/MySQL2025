from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation():
    # 初始化PPT
    prs = Presentation()

    # 定义一个辅助函数来创建带内容的幻灯片
    def add_slide(title, content_list, notes_text):
        slide_layout = prs.slide_layouts[1]  # 使用标题+内容布局
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title

        # 设置正文内容
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame

        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            if item.startswith("  -"):
                p.level = 1
            elif item.startswith("    -"):
                p.level = 2

        # 添加演讲稿 (备注)
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

        return slide

    # ================= 幻灯片 1: 封面 =================
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "基于Django的学生管理系统设计与实现"
    slide.placeholders[1].text = "汇报人：欣禹行\n学号：23012660"

    # 备注
    slide.notes_slide.notes_text_frame.text = (
        "各位老师、同学好。我是欣禹行。今天我汇报的主题是《基于Django的学生管理系统设计与实现》。"
        "本次实验我参考了机票预订系统的分层架构，并在此基础上自主设计并开发了一套完整的考勤管理模块。"
    )

    # ================= 幻灯片 2: 项目背景与开发环境 =================
    content = [
        "系统架构：采用Django多应用架构（Apps）",
        "  - 分离式设计：Student, Teacher, Admin 三大门户独立",
        "开发环境：",
        "  - 语言框架：Python 3.8+ / Django 5.0.7",
        "  - 数据库：MySQL 8.0 (使用原生SQL，非ORM)",
        "  - 前端：Bootstrap 5.3.3 + 原生JS (Fetch API)",
        "核心亮点：",
        "  - 自主设计考勤管理闭环（点名-请假-审批-自动化）"
    ]
    notes = (
        "首先介绍一下项目的整体情况。系统采用了类似机票预订系统的多应用架构，将学生、教师和管理员的逻辑完全分离，便于权限管理。"
        "技术栈方面，后端使用了Django 5.0.7，为了加深对数据库的理解，我使用了mysqlclient驱动进行原生SQL操作，而不是依赖ORM。"
        "前端使用了Bootstrap和AJAX技术。本系统的最大亮点是我自主设计了一套包含点名、请假、审批和自动记录生成的考勤闭环模块。"
    )
    add_slide("项目概况与开发环境", content, notes)

    # ================= 幻灯片 3: 系统功能结构 =================
    content = [
        "1. 学生门户 (Student Portal)",
        "  - 个人信息、选课、成绩查询",
        "  - 考勤查询、请假申请（含图片上传）",
        "2. 教师门户 (Teacher Portal)",
        "  - 课程管理、成绩录入（Excel导入）",
        "  - 交互式考勤点名、请假审批",
        "3. 管理员门户 (Admin Portal)",
        "  - 基础数据增删改查、认证授权管理"
    ]
    notes = (
        "系统功能结构主要分为三个部分。"
        "学生端除了基础的选课和成绩查询外，重点实现了带图片上传的请假申请功能。"
        "教师端支持Excel成绩导入，以及我设计的交互式考勤点名和审批功能。"
        "管理员端则负责基础数据的维护和权限控制。"
    )
    add_slide("系统功能结构", content, notes)

    # ================= 幻灯片 4: 数据库设计 (重点) =================
    content = [
        "核心表结构：8个基础表 + 1个考勤扩展表",
        "重点设计表：",
        "  - attendances (考勤表)：记录单次考勤状态",
        "  - leave_requests (请假申请表)：考勤模块专用",
        "逻辑关系：",
        "  - 教师 <1:N> 课程",
        "  - 学生 <1:N> 请假申请",
        "  - 请假申请 <1:N> 考勤记录 (审批通过后自动关联)"
    ]
    notes = (
        "数据库方面，系统包含8个核心表。为了实现考勤功能，我特别设计了`leave_requests`表。"
        "这里的一个关键逻辑是：一张请假申请单，在审批通过后，会对应生成多条考勤记录。"
        "比如学生请假3天，审批后，attendances表中会自动生成3条状态为'请假'的记录，并通过外键与申请单关联。"
    )
    add_slide("数据库设计与ER逻辑", content, notes)

    # ================= 幻灯片 5: 核心模块 - 考勤管理系统 =================
    # 过渡页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "核心模块展示：自主设计的考勤系统"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "该模块实现了从“点名”到“请假”再到“审批”的完整数据闭环。"

    # 添加一个占位符提示插入流程图
    p = tf.add_paragraph()
    p.text = "\n[建议在此处插入报告中的'系统业务流程图'截图]"
    p.font.italic = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 0, 0)

    notes = (
        "接下来重点汇报我自主设计的考勤管理模块。"
        "这也是本次实验中逻辑最复杂、技术含量最高的部分。它不仅仅是简单的增删改查，而是一个完整的数据流转闭环。"
    )
    slide.notes_slide.notes_text_frame.text = notes

    # ================= 幻灯片 6: 考勤功能 1 - 教师交互式点名 =================
    content = [
        "交互设计：",
        "  - 创新的头像点击交互：出勤→迟到→早退→缺勤→请假",
        "  - 支持一键全选/批量保存",
        "技术实现：",
        "  - 前端：JS维护状态数组，Fetch API发送批量JSON",
        "  - 后端：batch_attendance() 视图",
        "  - ID生成：ATT + 学号后6位 + 时间戳",
        "  - SQL优化：原生SQL批量插入"
    ]
    notes = (
        "首先是教师点名功能。为了提高点名效率，我没有使用传统的下拉框，而是设计了点击头像切换状态的交互方式。"
        "技术上，前端使用JavaScript维护一个状态数组，点击保存时，通过Fetch API将JSON数据发送给后端。"
        "后端`batch_attendance`函数接收数据后，会生成唯一的考勤ID，并执行原生SQL进行批量插入，保证了性能。"
    )
    add_slide("考勤模块：教师交互式点名", content, notes)

    # ================= 幻灯片 7: 考勤功能 2 - 学生请假申请 =================
    content = [
        "功能特性：",
        "  - 支持表单填写与病假条图片上传",
        "  - 图片本地即时预览 (FileReader API)",
        "技术实现：",
        "  - 后端视图：submit_leave_request()",
        "  - 文件处理：Django request.FILES 接收",
        "  - 存储路径：自动保存至 MEDIA_ROOT/medical_certificates/",
        "  - ID生成：LV + 学号后8位 + 时间戳"
    ]
    notes = (
        "第二部分是学生请假申请。这里我实现了文件上传功能，允许学生上传病假条。"
        "为了提升用户体验，我使用了HTML5的FileReader API，让学生在提交前能即时预览图片。"
        "后端接收到`request.FILES`后，会将图片保存到配置好的媒体目录下，并在数据库中存储路径。"
    )
    add_slide("考勤模块：学生请假申请", content, notes)

    # ================= 幻灯片 8: 考勤功能 3 - 审批与自动化逻辑 =================
    content = [
        "业务逻辑核心：",
        "  - 审批不仅仅是状态更新，更是数据生产",
        "自动化算法 (Approval Logic)：",
        "  1. 验证权限：确认操作者是否为该课程教师",
        "  2. 更新状态：leave_requests 表状态置为'通过'",
        "  3. 自动生成：遍历请假日期范围 (Start -> End)",
        "  4. 插入记录：在 attendances 表中自动插入'请假(5)'记录",
        "  5. 关联数据：将考勤记录与请假申请ID绑定"
    ]
    notes = (
        "最后是考勤模块最核心的审批逻辑。这也是本系统的难点所在。"
        "当教师点击“通过”时，后端不仅更新申请单的状态，还会触发一个自动化算法。"
        "系统会计算请假开始到结束的日期范围，自动为每一天生成一条状态为“请假”的考勤记录，并插入数据库。"
        "这样教师就不用再去手动修改考勤表，实现了真正的自动化管理。"
    )
    add_slide("考勤模块：审批与自动化逻辑", content, notes)

    # ================= 幻灯片 9: 总结与收获 =================
    content = [
        "实验总结：",
        "  - 完成了全栈Web开发流程，实现了多⻆色权限隔离",
        "  - 深入理解了Django的MVT模式与原生SQL的结合",
        "  - 掌握了文件上传、AJAX异步交互等前端技术",
        "不足与展望：",
        "  - 界面美观度有待提升",
        "  - 审批流程可以增加邮件通知功能"
    ]
    notes = (
        "总结本次实验，我成功开发了一个功能完备的学生管理系统。通过这个项目，我不仅掌握了Django框架的基础，"
        "更重要的是通过自主设计考勤模块，深入理解了前后端数据交互和复杂的业务逻辑处理。"
        "未来可以进一步优化界面UI，并增加邮件通知等功能。汇报完毕，谢谢大家。"
    )
    add_slide("总结与收获", content, notes)

    # 保存PPT
    output_file = "学生管理系统_实验汇报_欣禹行.pptx"
    prs.save(output_file)
    print(f"PPT已生成: {output_file}")
    print("请在打开PPT后，点击下方的'备注'按钮查看生成的演讲稿。")


if __name__ == "__main__":
    create_presentation()